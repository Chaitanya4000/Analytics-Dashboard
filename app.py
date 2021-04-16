# >>>>>>>>>>>>>>> Analytics Dashboard <<<<<<<<<<<<<<<

# >>>>>>>>>>>>>>> Import all required libraries <<<<<<<<<<<<<<<
from tkinter import *
from PIL import ImageTk, Image
import operator
import time
import os
from langdetect import detect
import re
from colorama import win32
from translate import Translator
import pandas as pd
import numpy as np
from tkinter.ttk import Progressbar
import matplotlib.pyplot as plt
from tkinter import filedialog, messagebox
from datetime import datetime
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_DATA_LABEL_POSITION
from pptx.util import Cm
import nltk
from nltk.corpus import stopwords
from wordcloud import WordCloud, STOPWORDS
from gensim.parsing.preprocessing import remove_stopwords
from tkinter import messagebox as mbox
import multiprocessing
#import pypiwin32.client as win32

# >>>>>>>>>>>>>>> Initialize Variables <<<<<<<<<<<<<<<
str_OutputFolderPath = ""
str_InputFilePath = ""
str_InputFileName = ""
str_OutputFileName = ""
str_inputsheetname = ""
Li_ColumnNames = []
varx = []
vary = []
dic_chkvar = {}
dic_dropdwnvar = {}
bol_loginsuccess=False


def KillMainApp():
    app.destroy()

class Login:
    def __init__(self,app):
        self.app=app
        self.app.title("Login System")
        self.app.geometry("750x600")
        self.app.resizable(False,False)

        #===BG Image=====
        self.bg = ImageTk.PhotoImage(file="images/BgImg.PNG")
        self.bg_image=Label(self.app, image=self.bg).place(x=0,y=0,relwidth=1,relheight=1)

        # ====Login Frame=====
        Frame_login=Frame(self.app,bg="white")
        Frame_login.place(x=130,y=130,height=340,width=500)

        #title of login page
        title = Label(Frame_login,text="Login Here",font=("Impact",25,"bold"),fg="#d77337",bg="white").place(x=90,y=30)

        lbl_user = Label(Frame_login, text="Username", font=("Goudy Old Style", 15, "bold"), fg="gray",
                     bg="white").place(x=90,
                                       y=140)
        self.txt_user=Entry(Frame_login,font=("times new roman", 15),bg="lightgray")
        self.txt_user.place(x=90,y=170,width=350,height=35)

        lbl_pass = Label(Frame_login, text="Password", font=("Goudy Old Style", 15, "bold"), fg="gray",
                         bg="white").place(x=90,
                                           y=210)
        self.txt_pass = Entry(Frame_login, font=("times new roman", 15), bg="lightgray")
        self.txt_pass.place(x=90, y=240, width=350, height=35)
        self.txt_pass.config(show="*")

        #forget_btn = Button(Frame_login, text="Forget Password?",cursor="hand2", bg="white",fg="#d77337",bd=0,font=("times new roman",12)).place(x=90,y=280)
        Login_btn = Button(self.app,command=self.login_function,cursor="hand2", text="Login", fg="white", bg="#d77337",
                        font=("times new roman",20)).place(x=300, y=470,width=180,height=40)


    def login_function(self):

        if self.txt_pass.get()=="" or self.txt_user.get()=="":
            messagebox.showerror("Error","All fields are required",parent=self.app)
        elif self.txt_pass.get()!="admin" or self.txt_user.get()!="admin" :
            messagebox.showerror("Error","Invalid Username/Password",parent=self.app)
        else:
            KillMainApp()
            Login
            global bol_loginsuccess
            bol_loginsuccess=True



            #messagebox.showinfo("Welcome",f"Welcome {self.txt_user.get()}\nYour Password: {self.txt_pass.get()}", parent=self.app)

app=Tk()
obj=Login(app)
app.mainloop()

# >>>>>>>>>>>>>>> All functions that are used in applications <<<<<<<<<<<<<<<

# KillMainApp function will kill the main app window


def KillMainApp():
    app_root.destroy()


#Detect Language Function:
def langdetect(dataframe,columnname,str_InputFilePath,str_inputsheetname,q):
    for index, row in dataframe.iterrows():
        str_inputstring = re.sub(r'[0-9]', '1', row[columnname])
        str_inputstring = re.sub(r"[A-Za-z0-9._%+-]+"r"@[A-Za-z0-9.-]+"r"\.[A-Za-z]{2,4}",'test@test.com', str_inputstring)
        print("Detected language is: " + str(detect(str_inputstring)))
        q.put (str(detect(str_inputstring)))
    print("Language Detected")


# GenerateFilters function is called by generate filters button from main app.
def GenerateFilters():
    global str_OutputFileName
    global str_inputsheetname
    global str_OutputFolderPath
    str_OutputFolderPath = str(OutputFolderPath.get())
    str_inputsheetname = str(Inputsheetname.get())
    str_OutputFileName = os.path.basename(str_OutputFolderPath)
    print("Output folder path is: " + str_OutputFolderPath)
    print("Output file name is: " + str_OutputFileName)
    KillMainApp()

    OpenFilterWindow()


# Function for opening the file
def file_opener():
    global str_InputFileName
    global str_InputFilePath
    str_InputFilePath = filedialog.askopenfilenames(
        parent=app_root,
        initialdir='/',
        initialfile='tmp',
        filetypes=[("Excel", "xlsx")])
    print(str_InputFilePath)
    str_InputFilePath = str(str_InputFilePath).replace("(", "").replace(")", "").replace("'", "").replace(",", "")
    str_InputFileName = os.path.basename(str_InputFilePath)

    if (str_InputFilePath != ""):
        app_canvas.create_text(270, 200, text=str_InputFileName, font=("comicssansns", 10, "bold"), fill='White')


# OpenFilterWindow function is called to display all column names and filters for graphs.

def OpenFilterWindow():

    def SendMail():
        outlook = win32.Dispatch('outlook.application')
        mail = outlook.CreateItem(0)
        mail.To = str_ReceiverMailID.get()
        mail.Subject = "Results generated by SmarTAMM Analytics Solution"
        mail.Body = "Hello, \n \n Good day! \n \n PFA presentation which is generated by SmarTAMM Analytics Solution.\n \n BR,\n -SmarTAMM Analytics Solution"
        attachment = str_OutputFolderPath
        mail.Attachments.Add(attachment)
        mail.Send()

    def BuildGraphs():


        print("Build graph function executing...")

        pptx = Presentation("Templates\Template.pptx")
        first_slide_layout = pptx.slide_layouts[0]
        slide = pptx.slides.add_slide(first_slide_layout)
        slide.shapes.title.text = "< < < Data Analysis > > >"
        slide.placeholders[1].text = " -Created by SmarTAMM Analytics Solution."

        data = pd.read_excel(str_InputFilePath, str_inputsheetname)
        print("Output Folder Path: " + str_OutputFolderPath)
        print("Output File Name: " + str_OutputFileName)

        #if (Var_EnhancedChartButton.get() == 1):

            #categoryValues = []
            #categorycolumn = str_Category.get()
            #seriescolumn = str_series.get()

            #df = (data[categorycolumn].value_counts())
            #categoryValues = df.index.to_list()

            #df = (data[seriescolumn].value_counts())
            #SeriesValuesName = df.index.to_list()

            #slide = pptx.slides.add_slide(pptx.slide_layouts[5])
            #chart_data = ChartData()
            #chart_data.categories = categoryValues

            #seriesValues = []
            #for item in categoryValues:
             #   for name in SeriesValuesName:
              #      helloo = item
               #     hi = name
                #    newdf = data.query("{0} == @helloo & {1} == @hi".format(categorycolumn,seriescolumn))
                 #   seriesValues.append(len(newdf.index))

                #chart_data.add_series('series', seriesValues)
            #x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
            #graph_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
            #slide.shapes.title.text = "Test Chart"
            #bargraph = graph_frame.chart
            #category_axis = bargraph.category_axis
            #category_axis.has_major_gridlines = True

        int_CounterForColumn = 0

        for column in Li_ColumnNames:
            prgbar['value'] += 20
            newWindow_canvas.update_idletasks()

            print("Checking for column number: "+str(int_CounterForColumn)+" "+str(column))
            if (int_CounterForColumn == 20):
                break
            int_CounterForColumn = int_CounterForColumn+1
            if (dic_chkvar["ChkVar_" + column].get() == 1):
                print("Selected Column name is : " + column)
                print("Selected graph value for respective column is : " + str(
                    dic_dropdwnvar["DropdwnVar_" + column].get()))

                if (dic_dropdwnvar["DropdwnVar_" + column].get() == "BarGraph"):

                    df = (data[column].value_counts())
                    UniqueValues = df.index.to_list()
                    UniqueValuesCount = df.to_list()

                    print(UniqueValues)
                    print(UniqueValuesCount)

                    slide = pptx.slides.add_slide(pptx.slide_layouts[5])
                    chart_data = ChartData()
                    chart_data.categories = UniqueValues
                    chart_data.add_series('Series 1', UniqueValuesCount)

                    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
                    graph_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
                    slide.shapes.title.text = str(dic_dropdwnvar["DropdwnVar_" + column].get()) + " for " + column
                    bargraph = graph_frame.chart
                    category_axis = bargraph.category_axis
                    category_axis.has_major_gridlines = True

                    #Key observations on basis of analysis
                    if (Var_InterpretationButton.get() == 1):
                        dic_Analysis = {}
                        int_DictCounter = 0
                        for value in UniqueValues:
                            dic_Analysis[value] = UniqueValuesCount[int_DictCounter]
                            int_DictCounter = int_DictCounter+1
                        int_maxcount = max(dic_Analysis, key=dic_Analysis.get)
                        int_mincount = min(dic_Analysis, key=dic_Analysis.get)
                        int_meancount = sum(dic_Analysis.values()) / len(dic_Analysis)

                        Li_LessCount = []
                        Li_MaxCount = []
                        for k,v in dic_Analysis.items():
                            if v < int_meancount:
                                Li_LessCount.append(k)
                            if v > int_meancount:
                                Li_MaxCount.append(k)
                        Li_MaxCount = [str(i) for i in Li_MaxCount]
                        Li_LessCount = [str(i) for i in Li_LessCount]
                        slide = pptx.slides.add_slide(pptx.slide_layouts[1])
                        slide.shapes.title.text = str("Key observations for " + column)
                        bulletpointbox =  slide.shapes
                        bulletpoint1 = bulletpointbox.placeholders[1]
                        bulletpoint1.text =  "Below is details of analysis on basis of provided data:"
                        bulletpoint3 = bulletpoint1.text_frame.add_paragraph()
                        bulletpoint3.text = "The computed mean for " + str(column) + " is " + str(int_meancount)
                        bulletpoint3.level = 1
                        bulletpoint2 = bulletpoint1.text_frame.add_paragraph()
                        bulletpoint2.text = str(column) + " " + str(int_maxcount) + " has the maxmimum count."
                        bulletpoint2.level = 1
                        if (len(Li_MaxCount) != 0):
                            bulletpoint5 = bulletpoint1.text_frame.add_paragraph()
                            str_brk = ", "
                            bulletpoint5.text = str_brk.join(Li_MaxCount) + " has count more than average count."
                            bulletpoint5.level = 1
                        bulletpoint4 = bulletpoint1.text_frame.add_paragraph()
                        bulletpoint4.text = str(column) + " " + str(int_mincount) + " has the lowest number."
                        bulletpoint4.level = 1
                        if (len(Li_LessCount) != 0):
                            bulletpoint5 = bulletpoint1.text_frame.add_paragraph()
                            str_brk = ", "
                            bulletpoint5.text = str(str_brk.join(Li_LessCount))+" has count less than average count."
                            bulletpoint5.level = 1


                if (dic_dropdwnvar["DropdwnVar_" + column].get() == "Language Translation"):
                    print("Executing language translation algorithm")
                    Li_TranslatedText = []
                    translator = Translator(to_lang="English")
                    for index, row in data.iterrows():
                        # Find all numbers and emailIds and replace them
                        str_inputstring = re.sub(r'[0-9]', '111', row[column])
                        str_inputstring = re.sub(r"[A-Za-z0-9._%+-]+"r"@[A-Za-z0-9.-]+"r"\.[A-Za-z]{2,4}",
                                                 'test@test.com', str_inputstring)
                        str_translatedtext = translator.translate(str_inputstring)
                        print("Translated text is: " + str_translatedtext)
                        Li_TranslatedText.append(str_translatedtext)
                    print(Li_TranslatedText)
                    data['Language translated for ' + column] = Li_TranslatedText
                    pdwriter = pd.ExcelWriter(str_InputFilePath, engine='xlsxwriter')
                    data.to_excel(pdwriter, sheet_name=str_inputsheetname)
                    pdwriter.save()
                    print("Language translation algorithm execution completed")

                if (dic_dropdwnvar["DropdwnVar_" + column].get() == "Lang Detect"):
                    print("Executing language detection algorithm")
                    print("Output Sheet name is :"+ str_inputsheetname)

                    Li_DetectedLanguage = []
                    np.array_split(data, 4)

                    if __name__ == "__main__":

                        q = multiprocessing.Queue()

                        p1 = multiprocessing.Process(target=langdetect,
                                                     args=((np.array_split(data, 4)[0]), column, str_InputFilePath, str_inputsheetname, q))
                        p2 = multiprocessing.Process(target=langdetect,
                                                     args=((np.array_split(data, 4)[1]), column, str_InputFilePath, str_inputsheetname, q))
                        p3 = multiprocessing.Process(target=langdetect,
                                                     args=((np.array_split(data, 4)[2]), column, str_InputFilePath, str_inputsheetname, q))
                        p4 = multiprocessing.Process(target=langdetect,
                                                     args=((np.array_split(data, 4)[3]), column, str_InputFilePath, str_inputsheetname, q))

                        p1.start()
                        p2.start()
                        p3.start()
                        p4.start()

                        p1.join()
                        p2.join()
                        p3.join()
                        p4.join()

                        while q.empty() is False:
                            while q.qsize() > 0:
                                Li_DetectedLanguage.append(q.get())
                            print(Li_DetectedLanguage)

                        print("Done!")


                    data['Language Detected for ' + column] = Li_DetectedLanguage
                    print(str_InputFilePath)
                    pdwriter = pd.ExcelWriter("C:\\Users\\A115618799\\Desktop\\Chaitanya\\PythonProjects\\AnalyticsDashboard_v2\\Vitesco Data\\Test.xlsx", engine='xlsxwriter')
                    data.to_excel(pdwriter, sheet_name=str_inputsheetname)
                    pdwriter.save()
                    print("Language detection algorithm execution completed")

                if (dic_dropdwnvar["DropdwnVar_" + column].get() == "PieChart"):

                    df = (data[column].value_counts())
                    UniqueValues = df.index.to_list()
                    UniqueValuesCount = df.to_list()

                    print(UniqueValues)
                    print(UniqueValuesCount)

                    slide = pptx.slides.add_slide(pptx.slide_layouts[5])
                    slide.shapes.title.text = str(dic_dropdwnvar["DropdwnVar_" + column].get()) + " for " + column
                    chart_data = ChartData()
                    chart_data.categories = UniqueValues
                    chart_data.add_series('Series 1', UniqueValuesCount)

                    x, y, cx, cy = Inches(0.5), Inches(2), Inches(5), Inches(4)
                    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

                    chart.has_legend = True
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    chart.legend.include_in_layout = False

                    chart.plots[0].has_data_labels = True
                    data_labels = chart.plots[0].data_labels
                    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

                    UniqueValuesCount2 = []
                    int_TotalNumber = sum(UniqueValuesCount)
                    print(int_TotalNumber)
                    for num in UniqueValuesCount:
                        int_percent = (num / int_TotalNumber)
                        print("Percentage value is :"+str(int_percent))
                        UniqueValuesCount2.append(int_percent)


                    chart_data = ChartData()
                    chart_data.categories = UniqueValues
                    chart_data.add_series('Series 1', UniqueValuesCount2)


                    x, y, cx, cy = Inches(5), Inches(2), Inches(5), Inches(4)
                    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

                    chart.has_legend = True
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    chart.legend.include_in_layout = False

                    chart.plots[0].has_data_labels = True
                    data_labels = chart.plots[0].data_labels
                    data_labels.number_format = '0.0%'
                    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

                    if (Var_InterpretationButton.get() == 1):
                        dic_Analysis = {}
                        int_DictCounter = 0
                        for value in UniqueValues:
                            dic_Analysis[value] = UniqueValuesCount[int_DictCounter]
                            int_DictCounter = int_DictCounter+1
                        int_maxcount = max(dic_Analysis, key=dic_Analysis.get)
                        int_mincount = min(dic_Analysis, key=dic_Analysis.get)
                        int_meancount = sum(dic_Analysis.values()) / len(dic_Analysis)
                        Li_LessCount = []
                        Li_MaxCount = []
                        for k,v in dic_Analysis.items():
                            if v < int_meancount:
                                Li_LessCount.append(k)
                            if v > int_meancount:
                                Li_MaxCount.append(k)
                        Li_MaxCount = [str(i) for i in Li_MaxCount]
                        Li_LessCount = [str(i) for i in Li_LessCount]
                        print(Li_MaxCount)
                        print(Li_LessCount)
                        slide = pptx.slides.add_slide(pptx.slide_layouts[1])
                        slide.shapes.title.text = str("Key observations for " + column)
                        bulletpointbox =  slide.shapes
                        bulletpoint1 = bulletpointbox.placeholders[1]
                        bulletpoint1.text =  "Below is details of analysis on basis of provided data:"
                        bulletpoint3 = bulletpoint1.text_frame.add_paragraph()
                        bulletpoint3.text = "The computed mean for " + str(column) + " is " + str(int_meancount)
                        bulletpoint3.level = 1
                        bulletpoint2 = bulletpoint1.text_frame.add_paragraph()
                        bulletpoint2.text = str(column) + " " + str(int_maxcount) + " has the maxmimum count."
                        bulletpoint2.level = 1
                        if (len(Li_MaxCount) != 0):
                            bulletpoint5 = bulletpoint1.text_frame.add_paragraph()
                            str_brk = ", "
                            bulletpoint5.text = str(str_brk.join(Li_MaxCount)) + " has count more than average count."
                            bulletpoint5.level = 1
                        bulletpoint4 = bulletpoint1.text_frame.add_paragraph()
                        bulletpoint4.text = str(column) + " " + str(int_mincount) + " has the lowest number."
                        bulletpoint4.level = 1
                        if (len(Li_LessCount) != 0):
                            bulletpoint5 = bulletpoint1.text_frame.add_paragraph()
                            str_brk = ", "
                            bulletpoint5.text = str_brk.join(Li_LessCount)+" has count less than average count."
                            bulletpoint5.level = 1


                if (dic_dropdwnvar["DropdwnVar_" + column].get() == "Unigram WC"):

                    InputColumn = data[column].dropna()
                    str_InputText = " "
                    for row in InputColumn:
                        str_InputText = str_InputText + " " + str(row)
                    print(str_InputText)


                    ListOfwords = []
                    ListOfwords = str_WordstoExclude.get().split(",")
                    print(ListOfwords[0])
                    str_InputText = getattr(str_InputText, 'upper')()
                    for word in ListOfwords:
                        print("Removing keyword :" + word)

                        word = getattr(word, 'upper')()
                        str_InputText = str_InputText.replace(word, "")


                    mask = np.array(Image.open("Images\\cloud.png"))
                    stopwords = set(STOPWORDS)
                    wc = WordCloud(background_color="black", mask=mask, max_words=200, stopwords=stopwords,collocations=False)
                    wc.generate(str_InputText)
                    wc.to_file("Images\\wc.png")
                    slide = pptx.slides.add_slide(pptx.slide_layouts[5])
                    slide.shapes.title.text = str(dic_dropdwnvar["DropdwnVar_" + column].get()) + " for " + column
                    img_WordcloudImage = "Images\\wc.png"
                    from_left = Inches(0.1)
                    from_top = Inches(2)
                    add_picture = slide.shapes.add_picture(img_WordcloudImage, from_left, from_top)

                if (dic_dropdwnvar["DropdwnVar_" + column].get() == "Bigram WC"):

                    InputColumn = data[column].dropna()
                    str_InputText = " "
                    for row in InputColumn:
                        str_InputText = str_InputText + " " + str(row)

                    ListOfwords = []
                    ListOfwords = str_WordstoExclude.get().split(",")
                    print(ListOfwords[0])
                    str_InputText = getattr(str_InputText, 'upper')()
                    for word in ListOfwords:
                        print("Removing keyword :" + word)

                        word = getattr(word, 'upper')()
                        str_InputText = str_InputText.replace(word, "")

                    from nltk.corpus import stopwords
                    stop_words = set(stopwords.words('english'))
                    stop_words.update(")", ".", ",", "!", "'", "(", "-", "--", "&", ":", "/", "'\'", "n't", "*")
                    stop_words.add("--")
                    Li_Bigrams = []
                    nltk_tokens = nltk.word_tokenize(str_InputText)
                    filtered_sentence = [w for w in nltk_tokens if not w in stop_words]
                    Li_Bigrams = list(nltk.bigrams(filtered_sentence))
                    Li_Bigrams = ['_'.join(tups) for tups in Li_Bigrams]
                    str_BigramOutput = " ".join(Li_Bigrams)

                    mask = np.array(Image.open("Images\\cloud.png"))
                    stopwords = set(STOPWORDS)
                    wc = WordCloud(background_color="black", mask=mask, max_words=200, stopwords=stopwords,collocations=False)
                    wc.generate(str_BigramOutput)
                    wc.to_file("Images\\wc1.png")
                    slide = pptx.slides.add_slide(pptx.slide_layouts[5])
                    slide.shapes.title.text = str(dic_dropdwnvar["DropdwnVar_" + column].get()) + " for " + column
                    img_WordcloudImage = "Images\\wc1.png"
                    from_left = Inches(0.1)
                    from_top = Inches(2)
                    add_picture = slide.shapes.add_picture(img_WordcloudImage, from_left, from_top)


                if (dic_dropdwnvar["DropdwnVar_" + column].get() == "WordCount"):

                    InputColumn = data[column].dropna()
                    str_InputText = " "
                    for row in InputColumn:
                        str_InputText = str_InputText + " " + str(row)

                    str_filteredsentence = remove_stopwords(str_InputText)
                    print(str_filteredsentence)
                    counts = dict()
                    words = str_filteredsentence.split()


                    for word in words:
                        if (word != "-"):
                            if word in counts:
                                counts[word] += 1
                            else:
                                counts[word] = 1

                    sorted_d = dict(sorted(counts.items(), key=operator.itemgetter(1), reverse=True))
                    print(sorted_d)

                    left_table = Inches(1)
                    top_table = Inches(1.5)
                    width_table = Inches(7)
                    height_table = Inches(1)

                    slide = pptx.slides.add_slide(pptx.slide_layouts[5])
                    slide.shapes.title.text = str(dic_dropdwnvar["DropdwnVar_" + column].get()) + " for " + column
                    add_table_Slide = slide.shapes.add_table(14, 2, left_table, top_table, width_table, height_table)
                    table1 = add_table_Slide.table
                    int_cellcounter = 1

                    cell = table1.cell(0, 0)
                    cell.text = "Word"

                    cell = table1.cell(0, 1)
                    cell.text = "Count"

                    l1 = list(sorted_d.keys())
                    l2 = list(sorted_d.values())

                    for i in l1:

                        cell = table1.cell(int_cellcounter, 0)
                        cell.text = i

                        cell2 = table1.cell(int_cellcounter, 1)
                        cell2.text = str(l2[int_cellcounter - 1])

                        int_cellcounter = int_cellcounter + 1

                        if (int_cellcounter == 14):
                            break

        if (dic_dropdwnvar["DropdwnVar_" + column].get() != "Language Detection") and (dic_dropdwnvar["DropdwnVar_" + column].get() != "Language Translation"):
            pptx.save(str_OutputFolderPath)
            mbox.showinfo("info","The ppt is saved at path :"+str_OutputFolderPath, parent=newWindow)


    # Get list of all columns in excel
    global Li_ColumnNames
    print("ExcelFilePath: " + str_InputFilePath)
    str_ExcelFilepath = str_InputFilePath
    data = pd.read_excel(str_ExcelFilepath,str_inputsheetname)
    Li_ColumnNames = data.columns

    # Building a new window
    newWindow = Tk()

    # sets the geometry of toplevel
    newWindow.geometry("1350x750")
    newWindow.maxsize(1350, 750)
    newWindow.minsize(1350, 750)

    # Add title to the main window
    newWindow.title("SmarTAMM Analytics Solution by T-systems")
    newWindow.configure(background="black")

    # Add backgroundImage
    img_BgImg = ImageTk.PhotoImage(Image.open("Images\\.img\\BgImg.PNG"))
    newWindow_canvas = Canvas(newWindow, width=1080, height=2160)
    newWindow_canvas.pack(fill="both", expand=True)
    newWindow_canvas.create_image(0, 0, image=img_BgImg, anchor="nw")

    # build title and project name
    newWindow_canvas.create_text(680, 35, text="SmarTAMM Analytics Solution", font=("comicssansns", 20, "bold"), fill='White')
    LogoImg = ImageTk.PhotoImage(Image.open("Images\\.img\\TSysLogo.PNG"))
    newWindow_canvas.create_image(680, 100, image=LogoImg)

    #Add progress bar to window
    prgbar = Progressbar(newWindow, orient=HORIZONTAL, length=300 , mode="determinate")
    Win_prgbar = newWindow_canvas.create_window(1170, 100, window=prgbar)

    newWindow_canvas.create_text(1170, 120, text="Progress Tracker", font=("comicssansns", 10, "bold"),
                                 fill='White')


    # Add separator line
    newWindow_canvas.create_line(1, 150, 2000, 150, fill="#fb0")

    # build label Column Name
    newWindow_canvas.create_text(740, 170, text="Select Graph Type:", font=("comicssansns", 10, "bold"),
                                 fill='White')

    # build label Column Name
    newWindow_canvas.create_text(540, 170, text="Select Column Names:", font=("comicssansns", 10, "bold"), fill='White')

    # build Exclude keywords from WordCloud input box
    newWindow_canvas.create_text(210, 170, text="Enter comma separated words to be excluded in WordCloud:", font=("comicssansns", 10, "bold"), fill='White')
    str_WordstoExclude = Entry(newWindow, width=50)
    newWindow_canvas.create_window(170, 200, window=str_WordstoExclude)

    #Interpretation based on charts
    Var_InterpretationButton = IntVar()
    Chk_InterpretationButton = Checkbutton(newWindow_canvas,variable=Var_InterpretationButton, text="Enable Interpretation", width=20)
    newWindow_canvas.create_window(100, 250, window=Chk_InterpretationButton)

    # build category input box
    #newWindow_canvas.create_text(155, 300, text="Enter category column for enhanced charts:",font=("comicssansns", 10, "bold"), fill='White')
    #str_Category = Entry(newWindow, width=50)
    #newWindow_canvas.create_window(170, 330, window=str_Category)

    # build series input box
    #newWindow_canvas.create_text(150, 360, text="Enter series column for enhanced charts:",font=("comicssansns", 10, "bold"), fill='White')
    #str_series = Entry(newWindow, width=50)
    #newWindow_canvas.create_window(170, 390, window=str_series)

    #Enhanced Chart
    #Var_EnhancedChartButton = IntVar()
    #Chk_EnhancedChartButton = Checkbutton(newWindow_canvas,variable=Var_EnhancedChartButton, text="Create Enhanced Chart", width=20)
    #newWindow_canvas.create_window(100, 440, window=Chk_EnhancedChartButton)

    # Submit button on new window
    Btn_Submit = Button(newWindow, text='Generate output >>>', command=BuildGraphs)
    Win_Btn_Submit = newWindow_canvas.create_window(80, 330, window=Btn_Submit)

    # Add separator line
    newWindow_canvas.create_line(1, 510, 430, 510, fill="#fb0")

    # Add separator line
    newWindow_canvas.create_line(430,150,430,2000, fill="#fb0")

    # build Enter from mailID
    newWindow_canvas.create_text(90, 550, text="Enter receiver's mailID:", font=("comicssansns", 10, "bold"), fill='White')
    str_ReceiverMailID = Entry(newWindow, width=50)
    newWindow_canvas.create_window(170, 580, window=str_ReceiverMailID)

    #Send output mail button on new window
    Btn_SendMail = Button(newWindow, text='Send Output Via Mail >>>', command=SendMail)
    Win_Btn_Submit = newWindow_canvas.create_window(95, 630, window=Btn_SendMail)

    # Logic to display columns as a filter options.

    x = 1
    int_Rowcounter = 220
    int_ColumnCounter = 550
    int_Dropdowncounter = 0
    BGOptions1 = ["PieChart", "BarGraph", "Unigram WC", "Bigram WC", "WordCount"]

    for i in Li_ColumnNames:

        dic_chkvar["ChkVar_{0}".format(i)] = IntVar()
        dic_dropdwnvar["DropdwnVar_{0}".format(i)] = StringVar(newWindow)
        vary = list(dic_dropdwnvar)
        varx = list(dic_chkvar)
        print(varx[x - 1])
        print(vary[x - 1])

        Chk_Button = Checkbutton(newWindow_canvas, text=i, variable=dic_chkvar[varx[x - 1]], width=20)
        newWindow_canvas.create_window(int_ColumnCounter, int_Rowcounter, window=Chk_Button)
        dic_dropdwnvar[vary[x - 1]].set(BGOptions1[0])

        BGOption1 = OptionMenu(newWindow, dic_dropdwnvar[vary[x - 1]], *BGOptions1)
        BGOption1.config(width=8, font=('Helvetica', 10))
        newWindow_canvas.create_window(int_ColumnCounter + 180, int_Rowcounter, window=BGOption1)

        x = x + 1
        int_Rowcounter = int_Rowcounter + 50
        int_Dropdowncounter = int_Dropdowncounter + 50
        if (x == 11):
            # build label Column Name
            newWindow_canvas.create_text(1130, 170, text="Select Graph Type:",
                                         font=("comicssansns", 10, "bold"),
                                         fill='White')

            # build label Column Name
            newWindow_canvas.create_text(950, 170, text="Select Column Names:", font=("comicssansns", 10, "bold"),
                                         fill='White')

            int_ColumnCounter = int_ColumnCounter + 400
            int_Rowcounter = 220

        if (x == 22):
            break

    newWindow.mainloop()


# >>>>>>>>>>>>>>>>>>>> Build Main Application <<<<<<<<<<<<<<<<<<<<<<<<<<
if(bol_loginsuccess==True):
    app_root = Tk()

    # GUI Framework
    app_root.geometry("500x400")
    app_root.maxsize(500, 500)
    app_root.minsize(500, 500)

    app_root.title("SmarTAMM Analytics Solution by T-systems")
    app_root.configure(background="black")

    # Add backgroundImage
    image1 = ImageTk.PhotoImage(Image.open("Images\\.img\\BgImg.PNG"))
    app_canvas = Canvas(app_root, width=1080, height=2160)
    app_canvas.pack(fill="both", expand=True)
    app_canvas.create_image(0, 0, image=image1, anchor="nw")

    # build title and project name
    app_canvas.create_text(250, 35, text="SmarTAMM Analytics Solution", font=("comicssansns", 20, "bold"), fill='White')
    LogoImg = ImageTk.PhotoImage(Image.open("Images\\.img\\TSysLogo.PNG"))
    app_canvas.create_image(250, 100, image=LogoImg)

    # Add separator line
    app_canvas.create_line(1, 150, 1360, 150, fill="#fb0")

    # build input filepath input box
    app_canvas.create_text(85, 170, text="Enter Input File Path:", font=("comicssansns", 10, "bold"), fill='White')

    # Browse Button label
    Btn_Browse = Button(app_root, text='Browse & Select file', command=lambda: file_opener())
    Win_BtnBrowseWindow = app_canvas.create_window(75, 200, window=Btn_Browse)

    # Add separator line
    app_canvas.create_line(1, 230, 1360, 230, fill="#fb0")

    # build output folder input box
    app_canvas.create_text(105, 250, text="Enter Input Worksheet Name:", font=("comicssansns", 10, "bold"), fill='White')
    Inputsheetname = Entry(app_root, width=50)
    app_canvas.create_window(170, 280, window=Inputsheetname)

    # Add separator line
    app_canvas.create_line(1, 300, 1360, 300, fill="#fb0")

    # build output folder input box
    app_canvas.create_text(105, 330, text="Enter Output Folder path:", font=("comicssansns", 10, "bold"), fill='White')
    OutputFolderPath = Entry(app_root, width=50)
    app_canvas.create_window(170, 360, window=OutputFolderPath)

    # Add Generate Filter Button
    Btn_GenerateFilters = Button(app_root, text='Generate Filters > > >', command=GenerateFilters)
    Win_BtnGenerateFiltersWindow = app_canvas.create_window(80, 430, window=Btn_GenerateFilters)

    # Add separator line
    app_canvas.create_line(1, 380, 1360, 380, fill="#fb0")

    #Add separator line
    app_canvas.create_line(1, 500, 1360, 500, fill="#fb0")

    app_root.mainloop()